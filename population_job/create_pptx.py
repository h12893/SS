import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import pandas as pd
from IPython.display import Image as IPImage
import numpy as np

# PowerPointプレゼンテーションを作成
prs = Presentation()

# スライドを追加する関数
def add_slide(title, content=None, image_path=None):
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # コンテンツ
    if content:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content

    # 画像
    if image_path and os.path.exists(image_path):
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        slide.shapes.add_picture(image_path, left, top, width, height)

# タイトルスライド
add_slide(
    "人口動態とIT業務数の関係性分析",
    "プレゼンテーション資料\n\n作成日: 2026年5月3日\n分析者: 人口動態分析チーム"
)

# プロジェクト概要
add_slide(
    "プロジェクト概要",
    "目的:\n- 人口の多寡や増減がIT関係の業務数にどのように影響しているかを調査\n- 地方でのデータ活用の可能性を検証\n\n前提仮説:\n1. 人口が多い都市部の方がIT企業や案件が多い\n2. 人口が少ない地方の方がIT技術導入による効率化の恩恵をより受けられる\n\n背景:\n- 厚生労働省のハローワークデータでは紹介件数の減少が顕著\n- 地方でのデータ活用事例（広島県離島、茨城県高校）が注目されている"
)

# データソースと推計モデル
add_slide(
    "データソースと推計モデル",
    "使用データ:\n- 人口動向調査（厚生労働省）: 都道府県別人口データ\n- 経済センサス（総務省・経済産業省）: 情報通信業従業者数\n- 労働力調査（総務省）: ITエンジニア就業者数\n- 法人番号データベース（国税庁）: 情報通信業企業数\n\n推計モデルの構造:\n3層構造で実測に近い推計を実現:\n1. 全国ITエンジニア数（実測） - 労働力調査\n2. 都道府県別IT産業規模（実測） - 経済センサス×人口分布\n3. 職種別構成比（実測） - JISA調査等\n\n課題と対応:\n- ハローワークデータの信頼性低下 → 複数データソースの組み合わせ推計\n- 産業分類の不一致 → 経済センサスを基準とした整合性チェック"
)

# 分析結果1: データの概要
# 人口推移グラフを作成
years = [2012, 2016, 2021]
population_trend = [127000000, 126000000, 125000000]

plt.figure(figsize=(10, 6))
plt.plot(years, population_trend, marker='o', linewidth=2)
plt.title('日本の総人口推移', fontsize=16)
plt.xlabel('年', fontsize=12)
plt.ylabel('人口（万人）', fontsize=12)
plt.grid(True, alpha=0.3)
plt.tight_layout()
plt.savefig('temp_population_trend.png', dpi=150, bbox_inches='tight')
plt.close()

add_slide(
    "分析結果1: データの概要",
    "人口動態データの傾向",
    'temp_population_trend.png'
)

# 分析結果2: ITエンジニア数の推計
add_slide(
    "分析結果2: ITエンジニア数の推計",
    "全国レベルでの推計結果",
    'python/output/img/line_plot_労働力調査（就業者数）_性計_年齢計.png'
)

# 分析結果3: 地域別比較
add_slide(
    "分析結果3: 地域別比較",
    "都道府県別のIT産業規模",
    'python/output/img/line_plot_経済センサス（就業者数）_性計_全国.png'
)

# 分析結果4: クラスタリング分析
add_slide(
    "分析結果4: クラスタリング分析",
    "IT産業の地域クラスタ",
    'python/output/img/scatter_plot_推計経済センサス（就業者数）クラスタリング_東京都以外_K=4.png'
)

# 分析結果5: 人口とIT業務数の相関
add_slide(
    "分析結果5: 人口とIT業務数の相関",
    "主成分分析による関係性分析",
    'python/output/img/scatter_plot_推計経済センサス（就業者数）の主成分分析.png'
)

# 考察と示唆
add_slide(
    "考察と示唆",
    "主要な発見:\n1. 都市部集中の確認: 人口密集地域でのIT業務数の偏在が確認された\n2. 地方の潜在力: 人口規模に比してIT活用が進む地域の存在\n3. データ活用の重要性: 複数データソースの統合推計の有効性\n\n政策示唆:\n- 地方でのIT人材育成プログラムの強化\n- データ活用を通じた地域振興策の検討\n- ハローワーク以外の就職支援システムの構築\n\n今後の課題:\n- より詳細な職種別分析\n- リアルタイムデータの活用\n- 国際比較の実施"
)

# 結論
add_slide(
    "結論",
    "まとめ:\n- 人口動態とIT業務数の関係性を多角的に分析\n- 推計モデルの有効性を確認\n- 地方データ活用の可能性を示唆\n\n次のステップ:\n- 詳細分析の継続\n- 政策提言の作成\n- 実証実験の実施\n\n---\n\nご清聴ありがとうございました。"
)

# PowerPointファイルを保存
prs.save('population_job_presentation.pptx')

# 一時ファイルを削除
if os.path.exists('temp_population_trend.png'):
    os.remove('temp_population_trend.png')

print("PowerPointプレゼンテーションが作成されました: population_job_presentation.pptx")